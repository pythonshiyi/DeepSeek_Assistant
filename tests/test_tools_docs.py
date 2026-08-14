# -*- coding: utf-8 -*-
"""新工具（文档处理/RSS/二维码/KV/媒体/WebDAV）自测用例：主路径 + 错误路径。"""
import json
import os
import shutil
import subprocess
import sys
import tempfile
import time
import unittest
from types import SimpleNamespace
from unittest import mock

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import deepseek_client as dc
import permissions


def _init_perm(tmp):
    ws = os.path.join(tmp, "ws")
    os.makedirs(ws, exist_ok=True)
    permissions.init(os.path.join(tmp, "perm.json"), ws)
    data = permissions.get_data()
    data["filesystem"]["blocked_dirs"] = [
        d for d in data["filesystem"]["blocked_dirs"] if "AppData" not in d
    ]
    permissions.set_full_auto(True)
    return ws


class _FakeResp:
    def __init__(self, status=200, text="", content=None, headers=None):
        self.status_code = status
        self.text = text
        self.content = content if content is not None else text.encode("utf-8")
        self.headers = headers or {}

    def raise_for_status(self):
        if self.status_code >= 400:
            raise AssertionError(f"HTTP {self.status_code}")


class TestPdfRoundtrip(unittest.TestCase):
    """T01+T02 闭环：pdf_create 生成（含中文/Markdown）→ pdf_extract 读回。"""

    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_pdf_")
        self.ws = _init_perm(self.tmp)

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_create_md_then_extract(self):
        md = (
            "# 测试标题\n\n"
            "这是一个包含中文的正文段落。\n\n"
            "```python\nprint('hello')\n```\n\n"
            "| 列A | 列B |\n|---|---|\n| 1 | 2 |\n"
        )
        pdf = os.path.join(self.ws, "doc.pdf")
        out = dc.pdf_create(content=md, output=pdf, title="测试文档")
        self.assertIn("已生成 PDF", out)
        self.assertTrue(os.path.exists(pdf))
        # 闭环读回：文本模式
        txt = dc.pdf_extract(pdf)
        self.assertIn("测试标题", txt)
        self.assertIn("中文", txt)
        self.assertIn("第1页", txt)
        # 元数据模式
        meta = dc.pdf_extract(pdf, mode="meta")
        self.assertIn("测试文档", meta)
        self.assertIn("页数", meta)
        # 表格模式（PDF 内表格依赖 find_tables，读取不崩溃即可）
        t = dc.pdf_extract(pdf, mode="table")
        self.assertTrue(isinstance(t, str))

    def test_create_from_source_path(self):
        md_path = os.path.join(self.ws, "src.md")
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("# 源文件标题\n\n正文来自文件。")
        out = dc.pdf_create(source_path=md_path, output=os.path.join(self.ws, "s.pdf"))
        self.assertIn("已生成", out)
        self.assertTrue(os.path.exists(os.path.join(self.ws, "s.pdf")))

    def test_content_and_source_conflict(self):
        out = dc.pdf_create(content="a", source_path="b.md", output=os.path.join(self.ws, "x.pdf"))
        self.assertIn("二选一", out)

    def test_missing_output(self):
        self.assertIn("output", dc.pdf_create(content="x"))

    def test_extract_nonexistent(self):
        out = dc.pdf_extract(os.path.join(self.ws, "nope.pdf"))
        self.assertIn("不存在", out)

    def test_extract_bad_file(self):
        bad = os.path.join(self.ws, "fake.pdf")
        with open(bad, "w", encoding="utf-8") as f:
            f.write("这不是 PDF")
        out = dc.pdf_extract(bad)
        self.assertIn("错误", out)

    def test_extract_page_range(self):
        # 每节加长内容强制分页（reportlab 短内容会全挤在第一页）
        sections = []
        for i in range(1, 6):
            sections.append("## 第 %d 节内容\n" % i + "\n".join(f"第 {i} 节第 {j} 行填充内容" for j in range(1, 60)))
        md = "\n\n".join(sections)
        pdf = os.path.join(self.ws, "pages.pdf")
        dc.pdf_create(content=md, output=pdf)
        full = dc.pdf_extract(pdf)
        total = len(full.split("--- 第")) - 1
        self.assertGreaterEqual(total, 3, "测试 PDF 应至少 3 页")
        # 只取第 2 页：输出只含第 2 页标记，且不含第 3 页内容
        txt = dc.pdf_extract(pdf, pages="2")
        self.assertIn("第2页", txt)
        self.assertNotIn("第3页", txt)
        self.assertNotIn("第1页", txt)

    def test_extract_bad_range(self):
        pdf = os.path.join(self.ws, "r.pdf")
        dc.pdf_create(content="x", output=pdf)
        out = dc.pdf_extract(pdf, pages="0-5")
        self.assertIn("范围", out)

    def test_extract_encrypted(self):
        try:
            import fitz
        except ImportError:
            self.skipTest("PyMuPDF 未安装")
        pdf = os.path.join(self.ws, "enc.pdf")
        doc = fitz.open()
        doc.new_page().insert_text((72, 72), "secret")
        doc.save(pdf, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw="o", user_pw="u")
        doc.close()
        out = dc.pdf_extract(pdf)
        self.assertIn("加密", out)


class TestDocxRead(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_docx_")
        self.ws = _init_perm(self.tmp)
        from docx import Document

        doc = Document()
        doc.add_heading("主标题", 0)
        doc.add_heading("一级标题", 1)
        doc.add_heading("二级标题", 2)
        doc.add_paragraph("正文段落内容")
        doc.add_paragraph("项目要点", style="List Bullet")
        doc.add_paragraph("长段落填充。" * 60)  # 保证超 max_chars 可截断
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "列A"
        table.cell(0, 1).text = "列B"
        table.cell(1, 0).text = "1"
        table.cell(1, 1).text = "2"
        doc.save(os.path.join(self.ws, "t.docx"))

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_read_structure(self):
        out = dc.docx_read(os.path.join(self.ws, "t.docx"))
        self.assertIn("# 主标题", out)
        self.assertIn("# 一级标题", out)
        self.assertIn("## 二级标题", out)
        self.assertIn("正文段落内容", out)
        self.assertIn("- 项目要点", out)
        self.assertIn("| 列A |", out)

    def test_old_doc_extension(self):
        p = os.path.join(self.ws, "old.doc")
        with open(p, "wb") as f:
            f.write(b"x")
        out = dc.docx_read(p)
        self.assertIn(".docx", out)

    def test_max_chars_truncation(self):
        out = dc.docx_read(os.path.join(self.ws, "t.docx"), max_chars=20)
        self.assertIn("已截断", out)
        self.assertLess(len(out), 1000)

    def test_missing_file(self):
        self.assertIn("不存在", dc.docx_read(os.path.join(self.ws, "n.docx")))


class TestPptxRead(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_pptx_")
        self.ws = _init_perm(self.tmp)
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "季度汇报"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = "要点一"
        p = tf.add_paragraph()
        p.text = "要点二"
        slide.notes_slide.notes_text_frame.text = "开场先讲背景"
        prs.save(os.path.join(self.ws, "t.pptx"))

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_read_with_notes(self):
        out = dc.pptx_read(os.path.join(self.ws, "t.pptx"))
        self.assertIn("幻灯片数: 1", out)
        self.assertIn("标题: 季度汇报", out)
        self.assertIn("要点一", out)
        self.assertIn("要点二", out)
        self.assertIn("备注: 开场先讲背景", out)

    def test_read_without_notes(self):
        out = dc.pptx_read(os.path.join(self.ws, "t.pptx"), include_notes=False)
        self.assertNotIn("备注", out)

    def test_old_ppt_extension(self):
        p = os.path.join(self.ws, "old.ppt")
        with open(p, "wb") as f:
            f.write(b"x")
        self.assertIn(".pptx", dc.pptx_read(p))


class TestRssFetch(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_rss_")
        self.ws = _init_perm(self.tmp)
        dc.RSS_SOURCES_FILE = os.path.join(self.tmp, "rss.json")

    def tearDown(self):
        dc.RSS_SOURCES_FILE = None
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_add_list_remove(self):
        out = dc.rss_fetch(action="add", url="https://example.com/feed.xml")
        self.assertIn("已添加", out)
        lst = dc.rss_fetch(action="list")
        self.assertIn("example.com", lst)
        out2 = dc.rss_fetch(action="remove", url="https://example.com/feed.xml")
        self.assertIn("已移除", out2)
        self.assertIn("没有 RSS 订阅", dc.rss_fetch(action="list"))
        # 持久化验证
        self.assertTrue(os.path.exists(dc.RSS_SOURCES_FILE))

    def test_add_duplicate(self):
        dc.rss_fetch(action="add", url="https://example.com/f.xml")
        self.assertIn("已订阅", dc.rss_fetch(action="add", url="https://example.com/f.xml"))

    def test_bad_action(self):
        self.assertIn("action", dc.rss_fetch(action="fly"))

    def test_fetch(self):
        import calendar

        # 条目发布时间 = 1 小时前（保证在 since_hours 窗口内）
        published = time.gmtime(time.time() - 3600)
        entry = SimpleNamespace(
            title="AI 最新进展",
            link="https://example.com/ai",
            published_parsed=published,
            summary="<p>这是摘要内容</p>",
        )
        feed = SimpleNamespace(entries=[entry], feed=SimpleNamespace(title="科技博客"), bozo=0)
        with mock.patch.dict(sys.modules, {"feedparser": mock.MagicMock(parse=staticmethod(lambda *a, **k: feed))}):
            out = dc.rss_fetch(action="fetch", url="https://example.com/f.xml", since_hours=24)
        self.assertIn("AI 最新进展", out)
        self.assertIn("科技博客", out)

    def test_fetch_invalid_source(self):
        feed = SimpleNamespace(entries=[], bozo=1, bozo_exception="解析失败")
        with mock.patch.dict(sys.modules, {"feedparser": mock.MagicMock(parse=staticmethod(lambda *a, **k: feed))}):
            out = dc.rss_fetch(action="fetch", url="https://example.com/bad.xml")
        self.assertIn("无效的 RSS", out)


class TestQrcode(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_qr_")
        self.ws = _init_perm(self.tmp)

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_generate(self):
        p = os.path.join(self.ws, "q.png")
        out = dc.qrcode(action="generate", text="https://example.com", output=p)
        self.assertIn("已生成二维码", out)
        self.assertTrue(os.path.exists(p))
        with open(p, "rb") as f:
            self.assertEqual(f.read(8), b"\x89PNG\r\n\x1a\n")

    def test_generate_errors(self):
        self.assertIn("text", dc.qrcode(action="generate", text=""))
        self.assertIn("output", dc.qrcode(action="generate", text="x"))

    def test_read_without_pyzbar(self):
        # 模拟 pyzbar 缺失（CI 环境可能装有带 DLL 的 pyzbar wheel）→ 应降级提示而非崩溃
        p = os.path.join(self.ws, "q.png")
        dc.qrcode(action="generate", text="x", output=p)
        with mock.patch.dict("sys.modules", {"pyzbar": None, "pyzbar.pyzbar": None}):
            out = dc.qrcode(action="read", image_path=p)
        self.assertIn("pyzbar", out)

    def test_read_missing_image(self):
        self.assertIn("不存在", dc.qrcode(action="read", image_path=os.path.join(self.ws, "n.png")))


class TestKvStore(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_kv_")
        dc.KV_CACHE_DIR = os.path.join(self.tmp, "kv")

    def tearDown(self):
        dc.KV_CACHE_DIR = None
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_set_get_delete(self):
        self.assertIn("已写入", dc.kv_store(action="set", key="a", value="1"))
        self.assertIn("a: 1", dc.kv_store(action="get", key="a"))
        self.assertIn("已删除", dc.kv_store(action="delete", key="a"))
        self.assertIn("不存在", dc.kv_store(action="get", key="a"))

    def test_keys_and_search(self):
        dc.kv_store(action="set", key="name", value="张三")
        dc.kv_store(action="set", key="city", value="北京")
        keys = dc.kv_store(action="keys")
        self.assertIn("name", keys)
        self.assertIn("city", keys)
        hit = dc.kv_store(action="search", pattern="张三")
        self.assertIn("name", hit)
        miss = dc.kv_store(action="search", pattern="不存在的词xyz")
        self.assertIn("未找到", miss)

    def test_ttl_expiry(self):
        dc.kv_store(action="set", key="t", value="v", ttl_seconds=1)
        self.assertIn("t: v", dc.kv_store(action="get", key="t"))
        time.sleep(1.6)
        self.assertIn("已过期", dc.kv_store(action="get", key="t"))

    def test_bad_action(self):
        self.assertIn("action", dc.kv_store(action="fly"))

    def test_empty_value_length_guard(self):
        out = dc.kv_store(action="set", key="k", value="x" * (1024 * 1024 + 10))
        self.assertIn("1MB", out)


class TestMediaFfmpeg(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_media_")
        self.ws = _init_perm(self.tmp)
        try:
            import imageio_ffmpeg
        except ImportError:
            self.skipTest("imageio-ffmpeg 未安装")
        ff = imageio_ffmpeg.get_ffmpeg_exe()
        self.video = os.path.join(self.ws, "test.mp4")
        subprocess.run(
            [ff, "-y",
             "-f", "lavfi", "-i", "testsrc=duration=2:size=320x240:rate=10",
             "-f", "lavfi", "-i", "sine=frequency=440:duration=2",
             "-map", "0:v", "-map", "1:a", "-c:v", "libx264", "-c:a", "aac",
             "-pix_fmt", "yuv420p", "-shortest", self.video],
            capture_output=True, timeout=90,
        )
        self.assertTrue(os.path.isfile(self.video))

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_info(self):
        out = dc.media_ffmpeg(action="info", input=self.video)
        self.assertIn("时长", out)
        self.assertIn("test.mp4", out)

    def test_thumbnail(self):
        png = os.path.join(self.ws, "thumb.png")
        out = dc.media_ffmpeg(action="thumbnail", input=self.video, output=png, time="00:00:01")
        self.assertIn("已截图", out)
        self.assertTrue(os.path.exists(png))

    def test_extract_audio(self):
        mp3 = os.path.join(self.ws, "audio.mp3")
        out = dc.media_ffmpeg(action="extract_audio", input=self.video, output=mp3, format="mp3")
        self.assertIn("提取音频", out)
        self.assertTrue(os.path.exists(mp3))

    def test_bad_format(self):
        out = dc.media_ffmpeg(action="transcode", input=self.video,
                              output=os.path.join(self.ws, "x.exe"), format="exe")
        self.assertIn("format", out)

    def test_bad_time(self):
        out = dc.media_ffmpeg(action="thumbnail", input=self.video,
                              output=os.path.join(self.ws, "t.png"), time="abc")
        self.assertIn("time", out)

    def test_missing_input(self):
        self.assertIn("input", dc.media_ffmpeg(action="info"))


_WEBDAV_XML = """<?xml version="1.0"?>
<d:multistatus xmlns:d="DAV:">
  <d:response><d:href>/dir/</d:href>
    <d:propstat><d:prop><d:resourcetype><d:collection/></d:resourcetype></d:prop>
    <d:status>HTTP/1.1 200 OK</d:status></d:propstat></d:response>
  <d:response><d:href>/dir/a.txt</d:href>
    <d:propstat><d:prop><d:displayname>a.txt</d:displayname>
    <d:getcontentlength>12</d:getcontentlength><d:getlastmodified>Mon, 10 Aug 2026 09:00:00 GMT</d:getlastmodified></d:prop>
    <d:status>HTTP/1.1 200 OK</d:status></d:propstat></d:response>
</d:multistatus>"""


class TestWebdav(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.mkdtemp(prefix="dsa_wd_")
        self.ws = _init_perm(self.tmp)
        self.cfg_path = os.path.join(self.tmp, "webdav.json")
        with open(self.cfg_path, "w", encoding="utf-8") as f:
            json.dump({"url": "https://dav.example.com", "username": "u", "password": "p"}, f)
        dc.WEBDAV_CONFIG_FILE = self.cfg_path

    def tearDown(self):
        dc.WEBDAV_CONFIG_FILE = None
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_no_config(self):
        dc.WEBDAV_CONFIG_FILE = os.path.join(self.tmp, "missing.json")
        out = dc.webdav(action="list")
        self.assertIn("未配置", out)

    def test_list(self):
        class C:
            def request(self, method, url, **kw):
                return _FakeResp(status=207, text=_WEBDAV_XML)

        with mock.patch("deepseek_client._http_client", return_value=C()):
            out = dc.webdav(action="list", remote_path="/dir/")
        self.assertIn("a.txt", out)
        self.assertIn("DIR", out)

    def test_upload_download(self):
        src = os.path.join(self.ws, "up.txt")
        with open(src, "w", encoding="utf-8") as f:
            f.write("hello webdav")

        class C:
            def request(self, method, url, **kw):
                if method == "PUT":
                    return _FakeResp(status=201)
                if method == "GET":
                    return _FakeResp(status=200, content=b"hello webdav")
                raise AssertionError(method)

        with mock.patch("deepseek_client._http_client", return_value=C()):
            up = dc.webdav(action="upload", remote_path="/up.txt", local_path=src)
            self.assertIn("已上传", up)
            dl = os.path.join(self.ws, "dl.txt")
            down = dc.webdav(action="download", remote_path="/up.txt", local_path=dl)
            self.assertIn("已下载", down)
            with open(dl, encoding="utf-8") as f:
                self.assertEqual(f.read(), "hello webdav")

    def test_delete(self):
        class C:
            def request(self, method, url, **kw):
                return _FakeResp(status=204)

        with mock.patch("deepseek_client._http_client", return_value=C()):
            out = dc.webdav(action="delete", remote_path="/a.txt")
        self.assertIn("已删除", out)

    def test_missing_local(self):
        out = dc.webdav(action="upload", remote_path="/x", local_path=os.path.join(self.ws, "n.txt"))
        self.assertIn("不存在", out)

    def test_bad_action(self):
        self.assertIn("action", dc.webdav(action="fly"))


if __name__ == "__main__":
    unittest.main()
