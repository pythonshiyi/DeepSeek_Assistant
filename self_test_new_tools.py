# -*- coding: utf-8 -*-
"""新工具真实环境核验脚本（非 mock）：真实文件、真实网络、真实 ffmpeg。

用法：python self_test_new_tools.py
退出码 0 = 全部通过。
"""
import os
import shutil
import subprocess
import sys
import tempfile
import time

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import deepseek_client as dc
import permissions

PASS, FAIL = 0, 0


def check(name, ok, detail=""):
    global PASS, FAIL
    if ok:
        PASS += 1
        print(f"  [PASS] {name}")
    else:
        FAIL += 1
        print(f"  [FAIL] {name} :: {detail}")


def main():
    tmp = tempfile.mkdtemp(prefix="dsa_real_")
    ws = os.path.join(tmp, "ws")
    os.makedirs(ws, exist_ok=True)
    permissions.init(os.path.join(tmp, "perm.json"), ws)
    data = permissions.get_data()
    data["filesystem"]["blocked_dirs"] = [
        d for d in data["filesystem"]["blocked_dirs"] if "AppData" not in d
    ]
    permissions.set_full_auto(True)
    dc.KV_CACHE_DIR = os.path.join(tmp, "kv")
    dc.RSS_SOURCES_FILE = os.path.join(tmp, "rss.json")
    dc.WEBDAV_CONFIG_FILE = os.path.join(tmp, "webdav.json")

    print("== 1. PDF 闭环（生成含中文/Markdown → 提取文本/表格/元数据）==")
    md = (
        "# 核验报告\n\n"
        "本次核验包含中文内容：报告文档与数据表格。\n\n"
        "```python\nprint('hello whaletalk')\n```\n\n"
        "| 项目 | 数量 |\n|---|---|\n| 工具 | 9 |\n| 测试 | 249 |\n"
    )
    pdf = os.path.join(ws, "report.pdf")
    r = dc.pdf_create(content=md, output=pdf, title="核验报告")
    check("pdf_create 生成", "已生成" in r and os.path.exists(pdf), r[:100])
    size = os.path.getsize(pdf) if os.path.exists(pdf) else 0
    check("PDF 文件非空", size > 500, f"{size} 字节")
    txt = dc.pdf_extract(pdf)
    check("pdf_extract 文本（含中文）", "核验报告" in txt and "中文" in txt, txt[:150].replace("\n", "|"))
    check("pdf_extract 代码块内容", "hello whaletalk" in txt)
    meta = dc.pdf_extract(pdf, mode="meta")
    check("pdf_extract 元数据（标题）", "核验报告" in meta, meta.replace("\n", "|")[:150])
    tmode = dc.pdf_extract(pdf, mode="table")
    check("pdf_extract 表格模式不崩溃", isinstance(tmode, str))
    one = dc.pdf_extract(pdf, pages="1")
    check("pdf_extract 页码", "第1页" in one and "第2页" not in one)

    print("== 2. docx_read（真实 Word 文档）==")
    try:
        from docx import Document

        doc = Document()
        doc.add_heading("核验文档", 0)
        doc.add_heading("第一节", 1)
        doc.add_paragraph("Word 正文段落，验证中文读取。")
        doc.add_paragraph("要点甲", style="List Bullet")
        t = doc.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "名称"
        t.cell(0, 1).text = "值"
        t.cell(1, 0).text = "CPU"
        t.cell(1, 1).text = "8 核"
        docx = os.path.join(ws, "t.docx")
        doc.save(docx)
        out = dc.docx_read(docx)
        check("docx_read 标题层级", "# 核验文档" in out and "# 第一节" in out)
        check("docx_read 中文正文", "Word 正文段落" in out)
        check("docx_read 列表", "- 要点甲" in out)
        check("docx_read 表格", "| 名称 |" in out and "8 核" in out)
    except Exception as e:
        check("docx_read", False, str(e))

    print("== 3. pptx_read（真实 PPT，含组合形状）==")
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "季度汇报"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = "要点一"
        p = tf.add_paragraph()
        p.text = "要点二"
        # 组合形状：两个文本框放入 group
        group = slide.shapes.add_group_shape()
        box1 = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.4))
        box1.text_frame.text = "组合内文本"
        slide.notes_slide.notes_text_frame.text = "开场讲背景"
        pptx = os.path.join(ws, "t.pptx")
        prs.save(pptx)
        out = dc.pptx_read(pptx)
        check("pptx_read 标题", "标题: 季度汇报" in out, out[:200].replace("\n", "|"))
        check("pptx_read 要点", "要点一" in out and "要点二" in out)
        check("pptx_read 组合形状文本", "组合内文本" in out)
        check("pptx_read 备注", "备注: 开场讲背景" in out)
        out2 = dc.pptx_read(pptx, include_notes=False)
        check("pptx_read 备注开关", "备注" not in out2)
    except Exception as e:
        check("pptx_read", False, str(e))

    print("== 4. rss_fetch（真实网络抓取）==")
    try:
        feed_url = "https://www.ruanyifeng.com/blog/atom.xml"
        r = dc.rss_fetch(action="add", url=feed_url)
        check("rss add", "已添加" in r, r[:80])
        r2 = dc.rss_fetch(action="fetch", url=feed_url, limit=5, since_hours=0)
        has_entry = "1." in r2 and "http" in r2
        check("rss fetch 真实条目", has_entry, r2[:200].replace("\n", "|"))
        if has_entry:
            check("rss fetch 标题/链接/摘要齐全", "|" in r2)
        r3 = dc.rss_fetch(action="remove", url=feed_url)
        check("rss remove", "已移除" in r3)
    except Exception as e:
        check("rss_fetch（网络）", False, str(e))

    print("== 5. qrcode（真实生成 + PNG 校验）==")
    try:
        qr = os.path.join(ws, "q.png")
        r = dc.qrcode(action="generate", text="https://example.com/whaletalk", output=qr)
        check("qrcode generate", "已生成二维码" in r and os.path.exists(qr), r[:80])
        with open(qr, "rb") as f:
            check("qrcode PNG 魔数", f.read(4) == b"\x89PNG")
        r2 = dc.qrcode(action="read", image_path=qr)
        check("qrcode read 降级提示（pyzbar 未装）", "pyzbar" in r2, r2[:80])
    except Exception as e:
        check("qrcode", False, str(e))

    print("== 6. kv_store（真实持久化 + TTL）==")
    try:
        check("kv set", "已写入" in dc.kv_store(action="set", key="项目", value="核验任务"))
        check("kv get", "核验任务" in dc.kv_store(action="get", key="项目"))
        check("kv search", "项目" in dc.kv_store(action="search", pattern="核验"))
        dc.kv_store(action="set", key="ttl", value="x", ttl_seconds=1)
        time.sleep(1.5)
        check("kv TTL 过期", "已过期" in dc.kv_store(action="get", key="ttl"))
        check("kv delete", "已删除" in dc.kv_store(action="delete", key="项目"))
    except Exception as e:
        check("kv_store", False, str(e))

    print("== 7. media_ffmpeg（真实 ffmpeg 全链路）==")
    try:
        import imageio_ffmpeg

        ff = imageio_ffmpeg.get_ffmpeg_exe()
        video = os.path.join(ws, "v.mp4")
        subprocess.run(
            [ff, "-y", "-f", "lavfi", "-i", "testsrc=duration=2:size=320x240:rate=10",
             "-f", "lavfi", "-i", "sine=frequency=440:duration=2",
             "-map", "0:v", "-map", "1:a", "-c:v", "libx264", "-c:a", "aac",
             "-pix_fmt", "yuv420p", "-shortest", video],
            capture_output=True, timeout=90,
        )
        info = dc.media_ffmpeg(action="info", input=video)
        check("media info 时长", "时长" in info, info.replace("\n", "|")[:150])
        thumb = os.path.join(ws, "t.png")
        r = dc.media_ffmpeg(action="thumbnail", input=video, output=thumb, time="00:00:01")
        check("media thumbnail", "已截图" in r and os.path.exists(thumb), r[:80])
        mp3 = os.path.join(ws, "a.mp3")
        r = dc.media_ffmpeg(action="extract_audio", input=video, output=mp3, format="mp3")
        check("media extract_audio", "提取音频" in r and os.path.exists(mp3), r[:80])
        r = dc.media_ffmpeg(action="transcode", input=video, output=os.path.join(ws, "o.webm"), format="webm")
        check("media transcode webm", "已转码" in r and os.path.exists(os.path.join(ws, "o.webm")), r[:80])
    except Exception as e:
        check("media_ffmpeg", False, str(e))

    print("== 8. webdav（配置缺失路径核验）==")
    r = dc.webdav(action="list")
    check("webdav 未配置提示", "未配置" in r, r[:80])
    dc.WEBDAV_CONFIG_FILE = os.path.join(tmp, "missing.json")
    r2 = dc.webdav(action="list")
    check("webdav 配置缺失提示", "未配置" in r2)

    print("== 9. 可选依赖降级（模拟缺库）==")
    try:
        import builtins

        real_import = builtins.__import__

        def fake_import(name, *a, **k):
            if name == "fitz":
                raise ImportError("blocked for test")
            return real_import(name, *a, **k)

        builtins.__import__ = fake_import
        try:
            r = dc.pdf_extract(os.path.join(ws, "report.pdf"))
            check("pdf_extract 缺库提示", "pip_install PyMuPDF" in r, r[:80])
        finally:
            builtins.__import__ = real_import
    except Exception as e:
        check("缺库降级", False, str(e))

    shutil.rmtree(tmp, ignore_errors=True)
    print(f"\n==== 核验结果：{PASS} 通过 / {FAIL} 失败 ====")
    return 0 if FAIL == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
